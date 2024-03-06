import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# Step 1: Create the bar graph
data = {'Category A': 10, 'Category B': 20, 'Category C': 15}
categories = list(data.keys())
values = list(data.values())

plt.bar(categories, values)
plt.xlabel('Categories')
plt.ylabel('Values')
plt.title('Bar Graph')

# Step 2: Export the graph as an image
plt.savefig('bar_graph.png')
plt.close()

# Step 3: Open the PowerPoint presentation
ppt = Presentation('presentation.pptx')

# Step 4: Find and delete the existing text box (if any)
for slide in ppt.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if 'Replace this text' in paragraph.text:
                    sp = shape
                    slide.shapes._spTree.remove(sp._element)

# Step 5: Insert the created bar graph image
img_path = 'bar_graph.png'
left_inch = Inches(1)
top_inch = Inches(1)
width_inch = Inches(5)
height_inch = Inches(3)

ppt.slides[0].shapes.add_picture(img_path, left_inch, top_inch, width_inch, height_inch)

# Save the modified presentation
ppt.save('modified_presentation.pptx')



import matplotlib.pyplot as plt
import numpy as np

# Data
categories = ['Category A', 'Category B', 'Category C', 'Category D']
values1 = [20, 35, 30, 25]
values2 = [25, 30, 35, 40]

# Set up the figure and axis
fig, ax = plt.subplots()

# Set width of bar
bar_width = 0.35

# Set position of bar on X axis
r1 = np.arange(len(values1))
r2 = [x + bar_width for x in r1]

# Make the plot
plt.bar(r1, values1, color='b', width=bar_width, edgecolor='grey', label='Group 1')
plt.bar(r2, values2, color='r', width=bar_width, edgecolor='grey', label='Group 2')

# Add xticks on the middle of the group bars
plt.xlabel('Categories', fontweight='bold')
plt.xticks([r + bar_width/2 for r in range(len(values1))], categories)

# Add ylabel
plt.ylabel('Values', fontweight='bold')

# Add legend
plt.legend()

# Show plot
plt.show()





import matplotlib.pyplot as plt

data = [[1, 2, 3],
        [4, 5, 6],
        [7, 8, 9]]

# Create a figure and axis
fig, ax = plt.subplots()

# Hide the axes
ax.axis('off')

# Create the table
table = ax.table(cellText=data, loc='center', cellLoc='center')

# Merge cells for the header
header_cells = table.get_celld()
header_cells[(0, 0)].set_text('Merged Header')
header_cells[(0, 1)].visible_edges = ''
header_cells[(0, 2)].visible_edges = ''

# Adjust the cell heights
for i in range(3):
    header_cells[(0, i)].set_height(0.1)

# Adjust font size
table.auto_set_font_size(False)
table.set_fontsize(14)

# Adjust cell heights
table.scale(1, 1.5)

plt.show()

